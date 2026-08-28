#!/usr/bin/env python3
"""ERP An Binh Foods — Professional kiosk presentation (v2).

Cover hero, 5 diagram images, 4 mock UI screenshots, logo on every slide,
rounded cards, accent bars, progress dots, auto-advance + kiosk loop.
Rebuild: python3 docs/build-presentation-pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os, math, tempfile

# ── paths ──
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) if "__file__" in globals() else "/Users/vunam/Downloads/koola/ERP/ERP_V6"
DOCS = os.path.join(ROOT, "docs") if os.path.isdir(os.path.join(ROOT, "docs")) else "/Users/vunam/Downloads/koola/ERP/ERP_V6/docs"
OUT  = os.path.join(DOCS, "ERP-AnBinhFoods-Presentation-2026-08-21.pptx")
LOGO_ABF = "/Users/vunam/Downloads/koola/ERP/ERP_V6/frontend/public/abf-logo.png"
LOGO     = "/Users/vunam/Downloads/koola/ERP/ERP_V6/frontend/public/logo.png"
ASSETS = os.path.join(tempfile.gettempdir(), "pptx-assets-pro2")
os.makedirs(ASSETS, exist_ok=True)

# ── palette ──
NAVY   = RGBColor(0x0F,0x23,0x3A); NAVY2 = RGBColor(0x1E,0x3A,0x5F)
BLUE   = RGBColor(0x25,0x63,0xEB); BLUE_LIGHT = RGBColor(0xDB,0xE9,0xFE)
BLUE_MID = RGBColor(0x93,0xC5,0xFD); CYAN = RGBColor(0x06,0xB6,0xD4)
TEAL   = RGBColor(0x14,0xB8,0xA6); AMBER = RGBColor(0xF5,0x9E,0x0B)
AMBER_LIGHT = RGBColor(0xFF,0xFB,0xEB); GREEN = RGBColor(0x16,0xA3,0x4A)
SLATE_50=RGBColor(0xF8,0xFA,0xFC); SLATE_100=RGBColor(0xF1,0xF5,0xF9)
SLATE_200=RGBColor(0xE2,0xE8,0xF0); SLATE_400=RGBColor(0x94,0xA3,0xB8)
SLATE_500=RGBColor(0x64,0x74,0x8B); SLATE_700=RGBColor(0x33,0x41,0x55)
SLATE_900=RGBColor(0x0F,0x17,0x2A); WHITE=RGBColor(0xFF,0xFF,0xFF)
W = Inches(13.33); H = Inches(7.5)
DURATIONS = [10000,8000,12000,12000,12000,15000,13000,12000,9000,8000,14000,13000,9000,9000,10000,10000,10000,12000]
NS = {"p":"http://schemas.openxmlformats.org/presentationml/2006/main"}

def set_bg(slide, color):
    bg = slide.background; f = bg.fill; f.solid(); f.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill=None, line=None, radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.line.fill.background()
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = Pt(1)
    if radius: shp.adjustments[0]=0.08
    return shp

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=SLATE_900, alignment=PP_ALIGN.LEFT, font_name="Inter", line_spacing=None, italic=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.text=text
    p.font.size=Pt(font_size); p.font.bold=bold; p.font.italic=italic
    p.font.color.rgb=color; p.font.name=font_name; p.alignment=alignment
    if line_spacing: p.line_spacing=Pt(line_spacing)
    return tx

def pill(slide, left, top, text, bg=BLUE, fg=WHITE, font_size=7, width=Inches(3.85)):
    s = add_rect(slide, left, top, width, Inches(0.32), fill=bg, radius=True)
    s.text_frame.word_wrap=True
    s.text_frame.paragraphs[0].text=text
    s.text_frame.paragraphs[0].font.size=Pt(font_size)
    s.text_frame.paragraphs[0].font.color.rgb=fg
    s.text_frame.paragraphs[0].font.bold=True
    s.text_frame.paragraphs[0].font.name="Inter"
    s.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
    s.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    return s

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path): return None
    return slide.shapes.add_picture(path, left, top, width=width, height=height)

def add_logo_footer(slide, show_text=True):
    # small abf logo + slide number area
    if os.path.exists(LOGO_ABF):
        try: slide.shapes.add_picture(LOGO_ABF, Inches(0.5), Inches(7.08), height=Inches(0.28))
        except: pass
    if show_text:
        add_text_box(slide, Inches(1.55), Inches(7.12), Inches(4), Inches(0.2),
                     "ERP An Binh Foods  •  Intelligent Manufacturing Platform", 6, False, SLATE_400)

# ── assets: generate with PIL + matplotlib ──
def gen_assets():
    from PIL import Image, ImageDraw, ImageFont
    import pathlib
    def font(size, bold=False):
        # Arial handles small sizes correctly; Helvetica.ttc fails below 8pt
        candidates = ["/System/Library/Fonts/Supplemental/Arial.ttf","/Library/Fonts/Arial.ttf","/System/Library/Fonts/Helvetica.ttc"] if int(size) < 8 else ["/System/Library/Fonts/Helvetica.ttc","/System/Library/Fonts/Supplemental/Arial.ttf"]
        for pp in candidates:
            if os.path.exists(pp):
                try: return ImageFont.truetype(pp, int(size))
                except: pass
        return ImageFont.load_default()

    # hero
    hero_path = os.path.join(ASSETS, "hero.jpg")
    if not os.path.exists(hero_path):
        w,h = 1600, 950
        im = Image.new("RGB",(w,h),(15,35,58))
        d = ImageDraw.Draw(im,"RGBA")
        for y in range(h):
            t=y/h; r=int(15*(1-t)+28*t); g=int(35*(1-t)+58*t); b=int(58*(1-t)+95*t)
            d.line([(0,y),(w,y)],fill=(r,g,b))
        for x in range(-h,w,92):
            d.line([(x,0),(x+h,h)],fill=(255,255,255,6),width=1)
        d.ellipse([1080,-140,1820,600],outline=(245,158,11,190),width=3)
        d.ellipse([1110,-110,1790,570],outline=(245,158,11,55),width=1)
        for x in range(60,900,150):
            for y in range(640,900,32):
                d.ellipse([x,y,x+3,y+3],fill=(255,255,255,16))
        # factory blocks subtle
        for i,(x,ww,hh) in enumerate([(1050,130,190),(1195,170,240),(1380,110,160)]):
            a=20+i*6
            d.rounded_rectangle([x,700-hh,x+ww,700],radius=10,fill=(255,255,255,a))
            for wx in range(x+14,x+ww-8,24):
                for wy in range(700-hh+20,700-14,24):
                    d.rectangle([wx,wy,wx+11,wy+13],fill=(245,158,11,50))
        im.save(hero_path, quality=92)
    # grades chart
    grades_path = os.path.join(ASSETS, "grades.png")
    if not os.path.exists(grades_path):
        try:
            import matplotlib; matplotlib.use("Agg")
            import matplotlib.pyplot as plt, matplotlib.ticker as mticker
            grades=["A","B","B-Dầu","C","Vụn","Phế","Ướt","Khác"]; vals=[28,22,18,12,8,5,4,3]
            cols=["#0F233A","#1E3A5F","#2563EB","#06B6D4","#14B8A6","#F59E0B","#EF4444","#94A3B8"]
            fig,ax=plt.subplots(figsize=(10.5,2.7),dpi=220)
            fig.patch.set_facecolor("white"); ax.set_facecolor("white")
            bars=ax.bar(grades,vals,color=cols,edgecolor="white",linewidth=1.3,width=0.58,zorder=3)
            for b,v in zip(bars,vals): ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.7,f"{v}%",ha="center",va="bottom",fontsize=9,fontweight="bold",color="#0F233A")
            ax.set_ylim(0,35); ax.set_ylabel("Tỉ lệ (%)",fontsize=7,color="#64748B")
            ax.tick_params(axis="x",labelsize=9,colors="#0F233A",pad=6); ax.tick_params(axis="y",labelsize=7,colors="#94A3B8")
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
            ax.spines["left"].set_color("#E2E8F0"); ax.spines["bottom"].set_color("#E2E8F0")
            ax.yaxis.set_major_locator(mticker.MultipleLocator(10)); ax.grid(axis="y",color="#F1F5F9",linewidth=0.9,zorder=0)
            plt.tight_layout(pad=0.4); plt.savefig(grades_path,bbox_inches="tight",facecolor="white"); plt.close()
        except Exception as e: print("grades chart skip:",e)
    # mfg pipeline
    mfg_path = os.path.join(ASSETS, "mfg-pipeline.png")
    if not os.path.exists(mfg_path):
        w,h=1280,210; im=Image.new("RGBA",(w,h),(255,255,255,0)); d=ImageDraw.Draw(im,"RGBA")
        steps=[("Material\nStandard",(37,99,235),"Định mức"),("Material\nEvaluation",(245,158,11),"Brix + trừ kho"),("System\nOperation",(20,184,166),"4 giai đoạn sấy"),("Finished\nProduct",(22,163,74),"8 grades")]
        x=8; bw=282
        for i,(lab,col,sub) in enumerate(steps):
            d.rounded_rectangle([x+3,20,x+bw+3,150],radius=18,fill=(15,35,58,16))
            d.rounded_rectangle([x,16,x+bw,146],radius=18,fill=col,outline=(255,255,255),width=2)
            d.rounded_rectangle([x+10,24,x+bw-10,40],radius=8,fill=(255,255,255,30))
            fl=lab.split("\n"); fb=font(18,True); fs=font(11)
            d.text((x+bw//2,62),fl[0],fill="white",anchor="mm",font=fb); d.text((x+bw//2,86),fl[1],fill="white",anchor="mm",font=fb)
            d.text((x+bw//2,118),sub,fill=(255,255,255,210),anchor="mm",font=fs)
            if i<3:
                ax=x+bw+5; mid=h//2-2
                d.polygon([(ax,mid),(ax+20,mid-13),(ax+20,mid+13)],fill=(100,116,139)); d.rectangle([ax-12,mid-2,ax+2,mid+2],fill=(100,116,139))
            x+=bw+34
        im.save(mfg_path)
    # ai pipeline
    ai_path = os.path.join(ASSETS, "ai-pipeline.png")
    if not os.path.exists(ai_path):
        w,h=1320,150; im=Image.new("RGBA",(w,h),(255,255,255,0)); d=ImageDraw.Draw(im,"RGBA")
        steps=[("14 docs",(37,99,235)),("Dense +\nBM25",(30,58,95)),("RRF\nk=60",(6,182,212)),("FlashRank",(20,184,166)),("DeepSeek",(15,35,58)),("Faithfulness",(245,158,11)),("Cache",(100,116,139)),("Response",(22,163,74))]
        bw,gap=144,16; x0=8
        for i,(lab,col) in enumerate(steps):
            x=x0+i*(bw+gap)
            d.rounded_rectangle([x+2,18,x+bw+2,118],radius=14,fill=(15,35,58,14))
            d.rounded_rectangle([x,14,x+bw,114],radius=14,fill=col)
            fb=font(11,True)
            ls=lab.split("\n")
            if len(ls)==1: d.text((x+bw//2,64),ls[0],fill="white",anchor="mm",font=fb)
            else: d.text((x+bw//2,54),ls[0],fill="white",anchor="mm",font=fb); d.text((x+bw//2,74),ls[1],fill="white",anchor="mm",font=fb)
            if i<7:
                ax=x+bw+2; mid=64
                d.polygon([(ax+2,mid),(ax+11,mid-6),(ax+11,mid+6)],fill=(148,163,184))
        im.save(ai_path)
    # face layers
    face_path = os.path.join(ASSETS, "face-layers.png")
    if not os.path.exists(face_path):
        w,h=1280,210; im=Image.new("RGBA",(w,h),(255,255,255,0)); d=ImageDraw.Draw(im,"RGBA")
        layers=[("01  DeepFace\n+ MiniFASNet","0.5/0.5",(15,35,58)),("02  Frame\nQuality","35–225 / blur>12",(20,184,166)),("03  LBP Texture\n128×128","avg ≥0.35",(245,158,11)),("04  Temporal","bbox shift",(6,182,212))]
        bw=290; x0=8
        for i,(lab,sub,col) in enumerate(layers):
            x=x0+i*(bw+18)
            d.rounded_rectangle([x+2,16,x+bw+2,142],radius=16,fill=(15,35,58,12))
            d.rounded_rectangle([x,12,x+bw,138],radius=16,fill="white",outline=col,width=2)
            d.rounded_rectangle([x,12,x+bw,26],radius=6,fill=col)
            cx,cy=x+bw//2,56; d.ellipse([cx-20,cy-20,cx+20,cy+20],fill=col); d.text((cx,cy),f"{i+1}",fill="white",anchor="mm",font=font(14,True))
            ls=lab.split("\n"); d.text((x+bw//2,86),ls[0],fill=(15,35,58),anchor="mm",font=font(11,True))
            if len(ls)>1: d.text((x+bw//2,102),ls[1],fill=(15,35,58),anchor="mm",font=font(11,True))
            d.text((x+bw//2,126),sub,fill=(100,116,139),anchor="mm",font=font(9))
        d.rounded_rectangle([8,162,w-8,196],radius=10,fill=(15,35,58))
        d.text((w//2,179),"Final gate:  0.5 × anti  +  0.2 × temporal  +  0.15 × quality  +  0.15 × lbp  ≥  0.68",fill="white",anchor="mm",font=font(11,True))
        im.save(face_path)
    # warehouse hierarchy
    wh_path = os.path.join(ASSETS, "warehouse.png")
    if not os.path.exists(wh_path):
        w,h=1280,150; im=Image.new("RGBA",(w,h),(255,255,255,0)); d=ImageDraw.Draw(im,"RGBA")
        steps=[("Warehouse",(15,35,58)),("Lot\nzone CAD",(30,58,95)),("Slot\nK1.1…K1.n",(37,99,235)),("LotProduct\nsoLuong/giaThanh",(20,184,166)),("BM01 / BM03\nFIFO + reorder",(245,158,11))]
        bw,gap=232,18; x0=8
        for i,(lab,col) in enumerate(steps):
            x=x0+i*(bw+gap)
            d.rounded_rectangle([x+2,14,x+bw+2,122],radius=16,fill=(15,35,58,12))
            d.rounded_rectangle([x,10,x+bw,118],radius=16,fill=col)
            ls=lab.split("\n"); d.text((x+bw//2,44),ls[0],fill="white",anchor="mm",font=font(12,True))
            if len(ls)>1: d.text((x+bw//2,66),ls[1],fill=(255,255,255,215),anchor="mm",font=font(10))
            if i==2: d.text((x+bw//2,96),"capacity guard",fill=(255,255,255,150),anchor="mm",font=font(8))
            if i==4: d.text((x+bw//2,96),"transaction",fill=(255,255,255,150),anchor="mm",font=font(8))
            if i<4:
                ax=x+bw+3; d.polygon([(ax,h//2),(ax+13,h//2-9),(ax+13,h//2+9)],fill=(100,116,139))
        im.save(wh_path)
    # tech stack image
    tech_path = os.path.join(ASSETS, "tech.png")
    if not os.path.exists(tech_path):
        w,h=1240,430; im=Image.new("RGBA",(w,h),(255,255,255,0)); d=ImageDraw.Draw(im,"RGBA")
        layers=[("EXPERIENCE","React 18 + TS + Vite 5 + TailwindCSS + TanStack Query + Recharts","SPA 46 pages, 70 hooks",(37,99,235)),("APPLICATION","Node 20 + Express 5 + Prisma 6 + Zod + Winston + ws + web-push","82 routes, RBAC/ABAC",(20,184,166)),("DATA","PostgreSQL 16 (3 schemas, 60+ models, CUID) + Redis 7","Multi-schema, advisory locks",(245,158,11)),("INTELLIGENCE","FastAPI + DeepSeek + ChromaDB + FlashRank + ArcFace","RAG 14 docs + 72 tools + 4-layer face",(6,182,212)),("INFRASTRUCTURE","Docker Compose + Nginx (TLS 1.2/1.3, http2) + backup 3 lớp","Single-VPS, 6-phase playbook",(51,65,85))]
        y=6; bh=74
        for name,stack,role,col in layers:
            d.rounded_rectangle([4,y+2,w-4,y+bh+2],radius=12,fill=(15,35,58,10))
            d.rounded_rectangle([0,y,w-8,y+bh],radius=12,fill="white",outline=(226,232,240),width=1)
            d.rounded_rectangle([0,y,7,y+bh],radius=3,fill=col)
            d.text((18,y+12),name,fill=col,font=font(10,True))
            d.text((18,y+30),stack,fill=(51,65,85),font=font(9.5))
            d.text((18,y+50),role,fill=(148,163,184),font=font(8.5))
            cx=w-36; cy=y+bh//2; d.ellipse([cx-16,cy-16,cx+16,cy+16],fill=col); d.text((cx,cy),"◆",fill="white",anchor="mm",font=font(10))
            y+=bh+10
        im.save(tech_path)
    # mock screenshots
    def make_mock_table(path, title, headers, rows, accent):
        if os.path.exists(path): return
        w,h=760,420; im=Image.new("RGB",(w,h),"white"); d=ImageDraw.Draw(im,"RGBA")
        d.rounded_rectangle([0,0,w-1,h-1],radius=14,outline=(226,232,240),width=1)
        # header bar
        d.rounded_rectangle([0,0,w-1,38],radius=14,fill=accent)
        d.rounded_rectangle([0,22,w-1,38],radius=0,fill=accent)
        d.text((14,19),title,fill="white",anchor="lm",font=font(11,True))
        d.text((w-14,19),"● ● ●",fill=(255,255,255,160),anchor="rm",font=font(8))
        # column headers
        y=54; cw=w//len(headers)
        for i,hh_ in enumerate(headers):
            x=i*cw
            d.rectangle([x+6,y,x+cw-4,y+26],fill=(248,250,252),outline=(226,232,240),width=1)
            d.text((x+cw//2,y+13),hh_,fill=(51,65,85),anchor="mm",font=font(8,True))
        y+=34
        for r, row in enumerate(rows):
            bg=(255,255,255) if r%2==0 else (248,250,252)
            for i, val in enumerate(row):
                x=i*cw
                d.rectangle([x+6,y,x+cw-4,y+30],fill=bg,outline=(241,245,249),width=1)
                col=(15,35,58) if i==0 else (51,65,85)
                # pill status in last col
                if i==len(row)-1 and val in ("Đạt","HOAT_DONG","Đã nhập","Đã duyệt"):
                    d.rounded_rectangle([x+12,y+6,x+cw-12,y+24],radius=8,fill=(220,252,231))
                    d.text((x+cw//2,y+15),val,fill=(22,101,52),anchor="mm",font=font(7,True))
                else:
                    d.text((x+cw//2,y+15),val,fill=col,anchor="mm",font=font(8))
            y+=32
        # pagination bar
        d.rectangle([0,h-30,w-1,h-1],fill=(248,250,252),outline=(226,232,240),width=1)
        d.text((w//2,h-15),"‹  1  2  3  …  8  ›",fill=(100,116,139),anchor="mm",font=font(8))
        im.save(path, quality=92)

    make_mock_table(os.path.join(ASSETS,"mock-eval.png"), "Đánh giá nguyên liệu  •  MaterialEvaluation", ["maChien","Brix","Nhiệt độ","Lot","Kết quả"], [["MC-260821-01","18.5","42°C","K1.3-042","Đạt"],["MC-260821-02","16.2","38°C","K1.4-017","Đạt"],["MC-260821-03","19.0","45°C","K1.2-089","Đạt"],["MC-260821-04","15.8","40°C","K1.1-102","Xem lại"],["MC-260821-05","17.4","41°C","K1.3-045","Đạt"]], (245,158,11))
    make_mock_table(os.path.join(ASSETS,"mock-warehouse.png"), "Kho kiện  •  Lot / Slot / LotProduct", ["Slot","Sản phẩm","Số lượng","Giá thành","Trạng thái"], [["K1.1","Mít sấy A","120 kg","42,000","Đã nhập"],["K1.2","Mít sấy B","85 kg","38,000","Đã nhập"],["K1.3","Chuối sấy","60 kg","35,000","Đã nhập"],["K1.4","Xoài sấy","45 kg","48,000","Đã nhập"],["K1.5","Mít vụn","30 kg","18,000","Chờ"]], (20,184,166))
    make_mock_table(os.path.join(ASSETS,"mock-machine.png"), "Thiết bị  •  MachineSystem / MaintenancePlan", ["Máy","Trạng thái","Bảo trì tới","Lỗi gần nhất","Phụ tùng"], [["Sấy-01","HOAT_DONG","15/09/2026","—","—"],["Sấy-02","BAO_TRI","—","Quạt yếu","Cánh quạt"],["Chiên-01","HOAT_DONG","20/09/2026","—","—"],["Đóng gói-02","HOAT_DONG","10/09/2026","Kẹt băng","Dây curoa"]], (100,116,139))
    # chat mock
    chat_path=os.path.join(ASSETS,"mock-chat.png")
    if not os.path.exists(chat_path):
        w,h=760,420; im=Image.new("RGB",(w,h),"white"); d=ImageDraw.Draw(im,"RGBA")
        d.rounded_rectangle([0,0,w-1,h-1],radius=16,outline=(226,232,240),width=1)
        d.rounded_rectangle([0,0,w-1,44],radius=16,fill=(15,35,58))
        d.rounded_rectangle([0,28,w-1,44],radius=0,fill=(15,35,58))
        d.ellipse([14,10,32,28],fill=(37,99,235)); d.text((23,19),"AI",fill="white",anchor="mm",font=font(7,True))
        d.text((42,10),"Trợ lý xưởng  •  AI Assistant",fill="white",anchor="lm",font=font(10,True))
        d.text((42,26),"Hiểu quy trình ERP  •  14 tài liệu  •  72 tools",fill=(148,163,184),anchor="lm",font=font(7))
        d.text((w-14,16),"●",fill=(34,197,94),anchor="rm",font=font(10))
        # user bubble
        d.rounded_rectangle([260,62,w-14,102],radius=14,fill=(37,99,235))
        d.text((w-28,82),"Quy trình tạo YCBG thế nào?",fill="white",anchor="rm",font=font(9))
        # ai bubble
        d.rounded_rectangle([14,116,520,230],radius=14,fill=(248,250,252),outline=(226,232,240),width=1)
        d.text((28,132),"Theo SOP Kho (doc 03-Common-Management):",fill=(15,35,58),anchor="lm",font=font(8,True))
        d.text((28,152),"1. Tạo SupplyRequest với items và lot gợi ý",fill=(51,65,85),anchor="lm",font=font(8))
        d.text((28,170),"2. Hệ thống check FIFO K1.1 → K1.n và capacity guard",fill=(51,65,85),anchor="lm",font=font(8))
        d.text((28,188),"3. Duyệt → sinh PurchaseRequest (MANUAL/SHORTAGE)",fill=(51,65,85),anchor="lm",font=font(8))
        d.text((28,210),"[Nguồn: docs/chatbot/03-Common-Management.md]",fill=(100,116,139),anchor="lm",font=font(7))
        # action card
        d.rounded_rectangle([14,244,520,310],radius=12,fill="white",outline=(37,99,235),width=1)
        d.text((28,260),"⚡ Agent có thể tạo YCBG ngay — bạn có muốn tiếp tục?",fill=(15,35,58),anchor="lm",font=font(8,True))
        d.rounded_rectangle([28,278,160,298],radius=8,fill=(37,99,235)); d.text((94,288),"Tạo YCBG",fill="white",anchor="mm",font=font(8,True))
        d.rounded_rectangle([170,278,260,298],radius=8,fill="white",outline=(226,232,240),width=1); d.text((215,288),"Để sau",fill=(51,65,85),anchor="mm",font=font(8))
        # input bar
        d.rounded_rectangle([14,h-44,w-14,h-14],radius=22,fill=(248,250,252),outline=(226,232,240),width=1)
        d.text((28,h-29),"Hỏi bằng tiếng Việt…",fill=(148,163,184),anchor="lm",font=font(9))
        d.ellipse([w-52,h-38,w-22,h-16],fill=(37,99,235)); d.text((w-37,h-27),"↑",fill="white",anchor="mm",font=font(12,True))
        im.save(chat_path, quality=92)
    # dashboard mock
    dash_path=os.path.join(ASSETS,"mock-dashboard.png")
    if not os.path.exists(dash_path):
        w,h=760,420; im=Image.new("RGB",(w,h),"white"); d=ImageDraw.Draw(im,"RGBA")
        d.rounded_rectangle([0,0,w-1,h-1],radius=16,outline=(226,232,240),width=1)
        d.rounded_rectangle([0,0,w-1,42],radius=16,fill=(15,35,58)); d.rounded_rectangle([0,26,w-1,42],radius=0,fill=(15,35,58))
        d.text((14,21),"Dashboard  •  Tổng quan vận hành",fill="white",anchor="lm",font=font(10,True))
        # KPI cards
        kpis=[("Đơn hàng","128","↑ 12%"),("Mẻ chiên","86","maChien"),("Tồn kho","1,420 kiện","FIFO"),("Nhân sự","247","100% chấm")]
        x=10
        for title,val,sub in kpis:
            d.rounded_rectangle([x,54,x+176,122],radius=12,fill="white",outline=(226,232,240),width=1)
            d.rounded_rectangle([x,54,x+176,62],radius=6,fill=(248,250,252))
            d.text((x+10,64),title,fill=(100,116,139),anchor="lm",font=font(7,True))
            d.text((x+10,88),val,fill=(15,35,58),anchor="lm",font=font(16,True))
            d.text((x+10,106),sub,fill=(22,163,74),anchor="lm",font=font(7))
            x+=186
        # yield mini chart
        d.rounded_rectangle([10,132,460,360],radius=12,fill="white",outline=(226,232,240),width=1)
        d.text((22,146),"Yield theo grade — 7 ngày qua",fill=(15,35,58),anchor="lm",font=font(9,True))
        # bars
        grades=["A","B","B-D","C","Vụn","Phế"]; vals=[28,22,18,12,8,5]; cols=[(15,35,58),(30,58,95),(37,99,235),(20,184,166),(245,158,11),(239,68,68)]
        bx=32; bw=54; gap=16; base=330; maxv=30
        for i,(g,v,c) in enumerate(zip(grades,vals,cols)):
            h_ = int(160 * v / maxv)
            x=bx+i*(bw+gap)
            d.rounded_rectangle([x,base-h_,x+bw,base],radius=6,fill=c)
            d.text((x+bw//2,base+12),g,fill=(100,116,139),anchor="mm",font=font(7,True))
            d.text((x+bw//2,base-h_-10),f"{v}%",fill=(15,35,58),anchor="mm",font=font(7,True))
        # right: alerts
        d.rounded_rectangle([474,132,w-10,242],radius=12,fill=(255,251,235),outline=(253,224,71),width=1)
        d.text((488,148),"Cảnh báo",fill=(146,64,14),anchor="lm",font=font(9,True))
        d.text((488,168),"• Slot K1.7 sắp đầy (92%)",fill=(120,53,15),anchor="lm",font=font(7))
        d.text((488,184),"• Máy Sấy-02 tới hạn bảo trì",fill=(120,53,15),anchor="lm",font=font(7))
        d.text((488,200),"• 3 YCBG chờ duyệt",fill=(120,53,15),anchor="lm",font=font(7))
        d.rounded_rectangle([474,252,w-10,360],radius=12,fill=(240,253,244),outline=(187,247,208),width=1)
        d.text((488,268),"Mẻ gần nhất",fill=(22,101,52),anchor="lm",font=font(9,True))
        d.text((488,288),"MC-260821-01  •  Brix 18.5  •  8 grades  •  Đạt",fill=(21,128,61),anchor="lm",font=font(7))
        d.text((488,306),"→ Đã nhập kho bulk (daNhapKho ✓)",fill=(100,116,139),anchor="lm",font=font(7))
        im.save(dash_path, quality=92)

# ── slide builders ──
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, NAVY)
    # hero image on right
    hero = os.path.join(ASSETS, "hero.jpg")
    if os.path.exists(hero):
        add_image(s, hero, Inches(7.55), Inches(0.45), width=Inches(5.3), height=Inches(5.55))
        # thin amber frame
        add_rect(s, Inches(7.55), Inches(0.45), Inches(5.3), Inches(5.55), fill=None, line=AMBER, radius=True)
    add_rect(s, Inches(0.55), Inches(0.55), Inches(0.9), Inches(0.06), fill=AMBER)
    add_text_box(s, Inches(0.55), Inches(0.78), Inches(4), Inches(0.28), "AN BINH FOODS  •  ERP V6  •  INTELLIGENT MANUFACTURING PLATFORM", 7, True, BLUE_MID)
    add_text_box(s, Inches(0.55), Inches(1.18), Inches(6.8), Inches(1.2), "Từ mẻ chiên tới\nthành phẩm —", 34, True, WHITE, line_spacing=36)
    add_text_box(s, Inches(0.55), Inches(2.42), Inches(6.8), Inches(0.55), "mọi thứ được kết nối", 34, True, AMBER)
    add_text_box(s, Inches(0.55), Inches(3.18), Inches(6.8), Inches(0.45), "Intelligent Manufacturing Platform  •  cho nhà sản xuất trái cây sấy", 10, False, BLUE_MID)
    add_text_box(s, Inches(0.55), Inches(3.70), Inches(6.8), Inches(0.6), "Truy vết mẻ chiên  •  Kho kiện FIFO  •  Bảo trì theo kế hoạch\nAI hiểu ngữ cảnh xưởng  •  Kiosk khuôn mặt tại xưởng", 8.5, False, SLATE_200, line_spacing=13)
    pill_s = add_rect(s, Inches(0.55), Inches(4.65), Inches(3.35), Inches(0.42), fill=WHITE, radius=True)
    pill_s.text_frame.word_wrap=True; p=pill_s.text_frame.paragraphs[0]
    p.text="Sấy thông minh, quản lý liền mạch"; p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=NAVY; p.font.name="Inter"; p.alignment=PP_ALIGN.CENTER
    pill_s.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    if os.path.exists(LOGO_ABF):
        try:
            pic = s.shapes.add_picture(LOGO_ABF, Inches(0.55), Inches(5.32), height=Inches(0.52))
            pic.line.fill.background()
        except: pass
    add_text_box(s, Inches(0.55), Inches(6.85), Inches(12), Inches(0.25), "www.anbinhfoods.vn  •  Gian hàng [số]  —  Gặp chúng tôi để xem kiosk và mẻ chiên trực tiếp", 7, False, SLATE_400)
    for i in range(18):
        c=AMBER if i==0 else SLATE_500
        add_rect(s, Inches(0.55+i*0.22), Inches(6.38), Inches(0.12), Inches(0.12), fill=c, radius=True)
    return s

def slide_agenda(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.38), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.53), Inches(4), Inches(0.22), "AGENDA  •  15 PHÚT  •  6 PHẦN", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.78), Inches(9), Inches(0.45), "Hôm nay chúng ta sẽ đi qua gì?", 26, True, NAVY)
    # illustrated agenda row
    agenda_img = os.path.join(ASSETS, "agenda-row.png")
    if os.path.exists(agenda_img):
        add_image(s, agenda_img, Inches(0.50), Inches(1.55), width=Inches(12.33))
    else:
        steps=[("01","Vấn đề","Vì sao Excel không đủ"),("02","Giải pháp","5 trụ cột"),("03","Dòng chảy","Từ báo giá tới thành phẩm"),("04","Manufacturing","Mẻ chiên & 8-grade yield  ★"),("05","AI + Face","Kiosk và trợ lý xưởng"),("06","Nền tảng","Công nghệ, kiến trúc, bảo mật")]
        x0=Inches(0.6)
        for i,(num,title,desc) in enumerate(steps):
            x=x0+i*Inches(2.05)
            add_rect(s, x, Inches(1.65), Inches(1.85), Inches(2.35), fill=SLATE_50, radius=True, line=SLATE_200)
            add_text_box(s, x+Inches(0.15), Inches(1.8), Inches(0.4), Inches(0.25), num, 7, True, BLUE)
            add_text_box(s, x+Inches(0.15), Inches(2.08), Inches(1.55), Inches(0.3), title, 11, True, NAVY)
            add_text_box(s, x+Inches(0.15), Inches(2.48), Inches(1.55), Inches(0.4), desc, 7.5, False, SLATE_500, line_spacing=10)
            c=add_rect(s, x+Inches(0.15), Inches(3.18), Inches(0.42), Inches(0.42), fill=BLUE_LIGHT, radius=True)
            c.text_frame.paragraphs[0].text=str(i+1); c.text_frame.paragraphs[0].font.size=Pt(9); c.text_frame.paragraphs[0].font.bold=True; c.text_frame.paragraphs[0].font.color.rgb=BLUE; c.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
            c.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    add_rect(s, Inches(0.6), Inches(3.35), Inches(12.13), Inches(0.02), fill=SLATE_200)
    add_text_box(s, Inches(0.6), Inches(3.52), Inches(12), Inches(0.28), "Trọng tâm: Manufacturing  ★  •  Demo kiosk + ChatWidget tại gian hàng sau trình bày  •  Q&A 5 phút  •  15 phút thuyết trình + 5 phút demo + 5 phút Q&A", 7, False, SLATE_400, PP_ALIGN.CENTER)
    # bottom illustration mock strip
    dash = os.path.join(ASSETS, "mock-dashboard.png")
    if os.path.exists(dash):
        add_image(s, dash, Inches(2.2), Inches(4.00), width=Inches(8.9), height=Inches(1.15))
        add_rect(s, Inches(2.2), Inches(4.00), Inches(8.9), Inches(1.15), fill=None, line=SLATE_200, radius=True)
    add_logo_footer(s); return s

def slide_problem(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, SLATE_50)
    add_rect(s, Inches(0.6), Inches(0.38), Inches(0.6), Inches(0.05), fill=AMBER)
    add_text_box(s, Inches(0.6), Inches(0.53), Inches(4), Inches(0.22), "VẤN ĐỀ", 7, True, AMBER)
    add_text_box(s, Inches(0.6), Inches(0.78), Inches(10), Inches(0.45), "Quản lý nhà máy sấy bằng Excel — đã đến lúc thay đổi", 20, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.28), Inches(10), Inches(0.28), "Khi mẻ chiên không có dữ liệu, mọi quyết định đều là ước đoán.", 9, False, SLATE_500, italic=True)
    # illustrated problem cards
    illus=os.path.join(ASSETS, "problem-illus.png")
    if os.path.exists(illus):
        add_image(s, illus, Inches(1.6), Inches(1.72), width=Inches(10.1))
        # captions below illustration
        captions=[
            ("Truy vết mẻ chiên thủ công","Mã mẻ ghi tay, truy nguyên lô lỗi mất ngày"),
            ("Tồn kho kiện không khớp","Kho là con số tổng, FIFO phụ thuộc trí nhớ"),
            ("Chấm công gian lận & tốn công","Chấm hộ, ghi tay, tổng hợp thủ công"),
        ]
        for i,(title,desc) in enumerate(captions):
            x=Inches(1.85)+i*Inches(3.38)
            add_text_box(s, x, Inches(3.55), Inches(2.9), Inches(0.32), title, 8, True, NAVY, PP_ALIGN.CENTER)
            add_text_box(s, x, Inches(3.88), Inches(2.9), Inches(0.32), desc, 7, False, SLATE_500, PP_ALIGN.CENTER, line_spacing=9)
    else:
        cols=[("Truy vết mẻ chiên\nthủ công","Mã mẻ ghi tay, không xuyên suốt từ nguyên liệu tới thành phẩm.\nTruy nguyên lô lỗi mất ngày lần sổ sách.","◉"),("Tồn kho kiện\nkhông khớp","Kho là con số tổng, không phải vị trí K1.1.\nFIFO phụ thuộc trí nhớ, kiểm kê lệch liên tục.","▦"),("Chấm công gian lận\n& tốn công","Chấm hộ, ghi tay, tổng hợp cuối tháng thủ công.\nQuản lý mất giờ đối chiếu.","◐")]
        for i,(title,desc,icon) in enumerate(cols):
            x=Inches(0.6)+i*Inches(4.15)
            add_rect(s, x, Inches(1.85), Inches(3.85), Inches(3.0), fill=WHITE, radius=True, line=SLATE_200)
            ic=add_rect(s, x+Inches(0.25), Inches(2.08), Inches(0.52), Inches(0.52), fill=RGBColor(0xFE,0xF3,0xC7), radius=True)
            ic.text_frame.paragraphs[0].text=icon; ic.text_frame.paragraphs[0].font.size=Pt(14); ic.text_frame.paragraphs[0].font.color.rgb=AMBER; ic.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER; ic.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
            add_text_box(s, x+Inches(0.25), Inches(2.72), Inches(3.35), Inches(0.55), title, 10, True, NAVY, line_spacing=13)
            add_text_box(s, x+Inches(0.25), Inches(3.38), Inches(3.35), Inches(1.05), desc, 7.5, False, SLATE_500, line_spacing=11)
    add_rect(s, Inches(0.6), Inches(4.45), Inches(12.13), Inches(0.42), fill=RGBColor(0xFF,0xFB,0xEB), radius=True, line=RGBColor(0xFD,0xE6,0x8A))
    add_text_box(s, Inches(0.75), Inches(4.55), Inches(11.8), Inches(0.22), "Nếu mẻ chiên không có dữ liệu xuyên suốt, mọi báo cáo sau đó đều thiếu gốc.", 8, True, RGBColor(0x92,0x4C,0x0E), PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_solution(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.38), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.53), Inches(4), Inches(0.22), "GIẢI PHÁP", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.78), Inches(10), Inches(0.45), "Một nền tảng — mọi vận hành được kết nối", 20, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.28), Inches(10), Inches(0.28), "Không ghép nhiều phần mềm rời rạc. Một dòng chảy dữ liệu từ báo giá tới thành phẩm và nhân sự.", 8, False, SLATE_500)
    pillars=[("01  Traceability","maChien xuyên 4 bảng\n06:30 boundary, 8-grade yield\n4-stage drying + MachineSystem",BLUE, "icon-trace.png"),("02  Connected Ops","QuotationCalculator → Order\n→ Warehouse FIFO → Production\n→ Quality → Invoice/Payroll",TEAL, "icon-flow.png"),("03  AI-Assisted","RAG hybrid + 72-tool ReAct\nagent hiểu department scope\nHỏi tiếng Việt, thao tác thực",AMBER, "icon-ai.png"),("04  Physical-Digital","Kiosk ArcFace 4-layer liveness\nAdaptive gallery + advisory-lock\nKhuôn mặt là thẻ chấm công",CYAN, "icon-face.png"),("05  Extensible","Modular monolith, multi-schema\nRBAC/ABAC, WS+Web Push\nDocker Compose, backup 3 lớp",SLATE_700, "icon-extensible.png")]
    for i,(title,desc,col,icon_name) in enumerate(pillars):
        x=Inches(0.6)+i*Inches(2.5)
        add_rect(s, x, Inches(1.85), Inches(2.3), Inches(2.95), fill=WHITE, radius=True, line=SLATE_200)
        add_rect(s, x, Inches(1.85), Inches(2.3), Inches(0.06), fill=col, radius=True)
        # icon circle
        icon_path=os.path.join(ASSETS, icon_name)
        if os.path.exists(icon_path):
            try:
                s.shapes.add_picture(icon_path, x+Inches(0.78), Inches(1.95), width=Inches(0.74), height=Inches(0.74))
            except: pass
            add_text_box(s, x+Inches(0.15), Inches(2.78), Inches(2.0), Inches(0.28), title, 7.5, True, col, PP_ALIGN.CENTER)
            add_text_box(s, x+Inches(0.15), Inches(3.10), Inches(2.0), Inches(1.45), desc, 7, False, SLATE_500, line_spacing=10, alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(s, x+Inches(0.15), Inches(1.98), Inches(2.0), Inches(0.28), title, 7.5, True, col)
            add_text_box(s, x+Inches(0.15), Inches(2.35), Inches(2.0), Inches(1.6), desc, 7, False, SLATE_500, line_spacing=10)
    add_text_box(s, Inches(0.6), Inches(5.12), Inches(12), Inches(0.22), "8 nhóm năng lực  •  Sales & Quotation  •  Supply & Purchase  •  Warehouse (kiện/FIFO)  •  Manufacturing  •  Quality  •  Equipment  •  HR & Payroll  •  AI & Face", 6.5, False, SLATE_400, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_connected(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.38), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.53), Inches(6), Inches(0.22), "DÒNG CHẢY DỮ LIỆU", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.78), Inches(12), Inches(0.45), "Từ báo giá tới thành phẩm — một dòng chảy dữ liệu", 18, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.23), Inches(12), Inches(0.22), "Mỗi nghiệp vụ là một mắt xích — dữ liệu đi tiếp, không nhập lại. Một mã maChien nối suốt.", 8, False, SLATE_500)
    add_rect(s, Inches(0.6), Inches(1.65), Inches(12.13), Inches(1.05), fill=RGBColor(0xEF,0xF6,0xFF), radius=True, line=RGBColor(0xBF,0xDB,0xFE))
    add_text_box(s, Inches(0.75), Inches(1.70), Inches(2), Inches(0.18), "SALES & SUPPLY", 6, True, BLUE)
    add_text_box(s, Inches(0.75), Inches(1.92), Inches(11.5), Inches(0.65), "InternationalCustomer  →  QuotationRequest  →  QuotationCalculator (MaterialStandard + GeneralCost + ExportCost → giá hòa vốn)  →  Quotation (priceLocked)  →  Order (7 SX + 3 TT)  →  TaxReport\nSupplyRequest (multi-item)  →  PurchaseRequest (MANUAL/SHORTAGE/REORDER)  →  Supplier", 7, False, SLATE_700, line_spacing=10)
    add_rect(s, Inches(0.6), Inches(2.9), Inches(12.13), Inches(1.08), fill=AMBER_LIGHT, radius=True, line=RGBColor(0xFD,0xE6,0x8A))
    add_text_box(s, Inches(0.75), Inches(2.95), Inches(2.2), Inches(0.18), "WAREHOUSE & PRODUCTION  ★", 6, True, AMBER)
    add_text_box(s, Inches(0.75), Inches(3.18), Inches(11.5), Inches(0.65), "Warehouse → Lot (zone CAD) → Slot (K1.1…) → LotProduct (soLuong, giaThanh)  →  BM01 / BM03  (multi-line, FIFO, reorder hook)\nMaterialStandard → MaterialEvaluation (maChien, Brix, lotProductId → trừ kho) → SystemOperation (4 giai đoạn, HOAT_DONG guard) → FinishedProduct (8 grades, daNhapKho) → QualityEvaluation → ProductionReport", 7, False, SLATE_700, line_spacing=10)
    add_rect(s, Inches(0.6), Inches(4.18), Inches(12.13), Inches(0.65), fill=RGBColor(0xF0,0xFD,0xF4), radius=True, line=RGBColor(0xBB,0xF7,0xD0))
    add_text_box(s, Inches(0.75), Inches(4.23), Inches(1), Inches(0.18), "NHÂN SỰ", 6, True, GREEN)
    add_text_box(s, Inches(0.75), Inches(4.43), Inches(11.5), Inches(0.28), "Employee → Attendance (shift provenance, kiosk/face) → Timesheet → Payroll (OT 1.5/2/3) → Evaluation (QUICK/FULL, appeal, peer)", 7, False, SLATE_700)
    for i,t in enumerate(["Giá hòa vốn có cơ sở — khóa giá trước khi gửi khách","Kho kiện first-class — FIFO K1.1→K1.n, capacity guard","maChien xuyên 4 bảng — boundary 06:30"]):
        pill(s, Inches(0.6+i*4.15), Inches(5.15), t, bg=SLATE_900 if i==2 else SLATE_100, fg=WHITE if i==2 else SLATE_700, font_size=6.5)
    add_logo_footer(s); return s

def slide_mfg1(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=AMBER)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(5), Inches(0.20), "MANUFACTURING  ★  HERO 1/2", 7, True, AMBER)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Mỗi mẻ chiên đều có câu chuyện dữ liệu", 20, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.22), "Từ Brix ngâm tới 4 giai đoạn sấy — truy vết một mã maChien xuyên suốt", 8, False, SLATE_500, italic=True)
    # pipeline image
    p = os.path.join(ASSETS, "mfg-pipeline.png")
    if os.path.exists(p):
        add_image(s, p, Inches(0.55), Inches(1.55), width=Inches(12.23))
    else:
        add_text_box(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5), "MaterialStandard → MaterialEvaluation → SystemOperation → FinishedProduct", 9, False, SLATE_500, PP_ALIGN.CENTER)
    # mock screenshot
    mock = os.path.join(ASSETS, "mock-eval.png")
    if os.path.exists(mock):
        add_image(s, mock, Inches(0.6), Inches(3.05), width=Inches(5.9))
        add_text_box(s, Inches(0.6), Inches(5.32), Inches(5.9), Inches(0.18), "▲  Mock UI — Đánh giá nguyên liệu (Brix, nhiệt độ, lot)  •  thay bằng screenshot thực khi có", 6, False, SLATE_400, PP_ALIGN.CENTER)
    # right proofs
    add_rect(s, Inches(6.85), Inches(3.05), Inches(0.02), Inches(2.45), fill=SLATE_200)
    add_text_box(s, Inches(7.05), Inches(3.05), Inches(5.6), Inches(0.22), "3 ĐIỂM CHỨNG MINH", 7, True, BLUE)
    for i,(t,d) in enumerate([("Trừ kho gắn mẻ","Trừ LotProduct cụ thể qua lotProductId trong transaction — biết chính xác kiện nào cho mẻ nào."),("4 giai đoạn có guard","Chỉ vận hành khi MachineSystem ở HOAT_DONG — không sấy trên máy đang bảo trì/hỏng."),("Ngày sản xuất 06:30","Production day boundary 06:30 theo ca thực tế — báo cáo khớp ca xưởng.")] ):
        y=Inches(3.38)+i*Inches(0.72)
        add_rect(s, Inches(7.05), y, Inches(5.6), Inches(0.60), fill=SLATE_50, radius=True, line=SLATE_200)
        add_text_box(s, Inches(7.20), y+Inches(0.07), Inches(5.3), Inches(0.20), "✓  "+t, 7.5, True, NAVY)
        add_text_box(s, Inches(7.20), y+Inches(0.26), Inches(5.3), Inches(0.28), d, 6.5, False, SLATE_500, line_spacing=9)
    add_rect(s, Inches(7.05), Inches(5.78), Inches(5.6), Inches(0.42), fill=AMBER_LIGHT, radius=True, line=RGBColor(0xFD,0xE6,0x8A))
    add_text_box(s, Inches(7.15), Inches(5.88), Inches(5.4), Inches(0.22), "Không có mẻ nào là “không rõ nguồn gốc”.", 7.5, True, RGBColor(0x92,0x4C,0x0E), PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_mfg2(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=AMBER)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(5), Inches(0.20), "MANUFACTURING  ★  HERO 2/2", 7, True, AMBER)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "8-grade yield — hiểu rõ từng phần của thành phẩm", 19, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.22), "Không chỉ “bao nhiêu kg” — mà là A/B/B-Dầu/C/Vụn/Phế/Ướt, với tỉ lệ tự động", 8, False, SLATE_500, italic=True)
    # real chart image
    chart = os.path.join(ASSETS, "grades.png")
    if os.path.exists(chart):
        add_rect(s, Inches(0.55), Inches(1.52), Inches(12.23), Inches(2.05), fill=WHITE, radius=True, line=SLATE_200)
        add_image(s, chart, Inches(0.70), Inches(1.62), width=Inches(11.9))
    else:
        # fallback bars
        grades=[("A",(15,35,58)),("B",(30,58,95)),("B-Dầu",(6,182,212)),("C",(20,184,166)),("Vụn",(245,158,11)),("Phế",(239,68,68)),("Ướt",(148,163,184)),("Khác",(226,232,240))]
        x0=Inches(0.6)
        for i,(g,col) in enumerate(grades):
            x=x0+i*Inches(1.58)
            h=Inches(1.0) if g in ("A","B") else Inches(0.7) if g=="B-Dầu" else Inches(0.5)
            add_rect(s, x, Inches(2.8)-h, Inches(1.38), h, fill=RGBColor(*col), radius=True)
    for i,t in enumerate(["Tỉ lệ tự động — tiLe = grade / tổng mẻ","Nhập kho bulk idempotent — daNhapKho guard","Báo cáo KH/TT theo grade — biết ngay lệch"]):
        pill(s, Inches(0.6+i*4.15), Inches(3.95), t, bg=SLATE_900 if i==0 else SLATE_100, fg=WHITE if i==0 else SLATE_700, font_size=6.5)
    # grade legend row
    add_text_box(s, Inches(0.6), Inches(4.45), Inches(12.13), Inches(0.22), "A: loại 1 đẹp nhất (giá cao nhất)  •  B-Dầu: đặc thù chiên chân không  •  Vụn: tận dụng / bán vụn  •  Phế: loại bỏ  •  Ướt: cần sấy lại  •  tiLe tính tự động", 6.5, False, SLATE_400, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_warehouse(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=TEAL)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(6), Inches(0.20), "SUPPLY & WAREHOUSE", 7, True, TEAL)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Kiện là first-class — tồn kho không còn là con số mơ hồ", 18, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.22), "Mỗi kiện có vị trí K1.1, có FIFO, có capacity guard — không phải ô tổng hợp", 8, False, SLATE_500, italic=True)
    wh = os.path.join(ASSETS, "warehouse.png")
    if os.path.exists(wh):
        add_image(s, wh, Inches(0.55), Inches(1.50), width=Inches(12.23))
    mock = os.path.join(ASSETS, "mock-warehouse.png")
    if os.path.exists(mock):
        add_image(s, mock, Inches(0.6), Inches(2.85), width=Inches(6.1))
        add_text_box(s, Inches(0.6), Inches(5.18), Inches(6.1), Inches(0.18), "▲  Mock UI — Kho kiện Lot/Slot/LotProduct  •  thay bằng screenshot thực khi có", 6, False, SLATE_400, PP_ALIGN.CENTER)
        # right proofs
        add_rect(s, Inches(7.0), Inches(2.85), Inches(0.02), Inches(2.33), fill=SLATE_200)
        for i,(t,d) in enumerate([("FIFO thực tế","Xuất K1.1 trước, đúng nhập trước xuất trước — giảm tồn đọng"),("Capacity guard","Slot có giới hạn, hệ thống chặn khi vượt — không vượt tải kệ"),("BM01/BM03 + reorder","Multi-line transaction, tự gợi ý đặt lại khi tồn thấp")] ):
            y=Inches(2.85)+i*Inches(0.82)
            add_rect(s, Inches(7.20), y, Inches(5.5), Inches(0.70), fill=WHITE, radius=True, line=SLATE_200)
            add_text_box(s, Inches(7.35), y+Inches(0.10), Inches(5.2), Inches(0.22), t, 8, True, NAVY)
            add_text_box(s, Inches(7.35), y+Inches(0.32), Inches(5.2), Inches(0.30), d, 7, False, SLATE_500, line_spacing=10)
    else:
        for i,(t,d) in enumerate([("FIFO thực tế","Xuất K1.1 trước"),("Capacity guard","Chặn khi vượt"),("BM01/BM03 + reorder","Multi-line + reorder hook")] ):
            x=Inches(0.6)+i*Inches(4.15); add_rect(s, x, Inches(2.85), Inches(3.85), Inches(1.0), fill=WHITE, radius=True, line=SLATE_200)
            add_text_box(s, x+Inches(0.2), Inches(2.95), Inches(3.45), Inches(0.22), t, 8, True, NAVY)
            add_text_box(s, x+Inches(0.2), Inches(3.25), Inches(3.45), Inches(0.4), d, 7, False, SLATE_500)
    add_logo_footer(s); return s

def slide_engineering(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=SLATE_500)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(4), Inches(0.20), "ENGINEERING", 7, True, SLATE_500)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(10), Inches(0.45), "Bảo trì không còn là phản ứng — là kế hoạch", 20, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(10), Inches(0.22), "Từ hệ thống máy tới kế hoạch năm, tới fault KB và phụ tùng", 8, False, SLATE_500, italic=True)
    mock = os.path.join(ASSETS, "mock-machine.png")
    if os.path.exists(mock):
        add_image(s, mock, Inches(0.6), Inches(1.58), width=Inches(6.2))
        add_text_box(s, Inches(0.6), Inches(3.90), Inches(6.2), Inches(0.18), "▲  Mock UI — MachineSystem & MaintenancePlan", 6, False, SLATE_400, PP_ALIGN.CENTER)
        add_rect(s, Inches(7.05), Inches(1.58), Inches(0.02), Inches(2.32), fill=SLATE_200)
        # right diagram
        add_rect(s, Inches(7.25), Inches(1.58), Inches(5.45), Inches(1.15), fill=SLATE_50, radius=True, line=SLATE_200)
        add_text_box(s, Inches(7.40), Inches(1.72), Inches(5.15), Inches(0.85), "MachineSystem  (HOAT_DONG / BAO_TRI / HONG)\n→  MaintenancePlan  (kế hoạch năm, template-driven)\n→  FaultRecord  (ghi lỗi, tra KB)  →  SparePart  (phụ tùng)", 7.5, True, NAVY, line_spacing=10, alignment=PP_ALIGN.CENTER)
        for i,(t,d) in enumerate([("Template-driven","Kế hoạch lặp theo năm\nkhông tạo tay từng lần"),("Fault KB","Lỗi tra từ KB tích lũy\nthợ mới cũng biết cách xử"),("Gắn phụ tùng","SparePart gắn fault & machine\nbiết đã thay gì")] ):
            y=Inches(2.95)+i*0 if i==0 else Inches(2.95)
            # we stack vertically? Actually 3 cards below
            pass
        # 3 small cards
        for i,(t,d) in enumerate([("Template-driven","Kế hoạch lặp theo năm"),("Fault KB","Tra KB tích lũy"),("Gắn phụ tùng","Biết đã thay gì")] ):
            x=Inches(7.25)+i*Inches(1.85)
            add_rect(s, x, Inches(2.95), Inches(1.72), Inches(0.95), fill=WHITE, radius=True, line=SLATE_200)
            add_text_box(s, x+Inches(0.12), Inches(3.05), Inches(1.48), Inches(0.22), t, 7, True, NAVY, PP_ALIGN.CENTER)
            add_text_box(s, x+Inches(0.12), Inches(3.32), Inches(1.48), Inches(0.40), d, 6.5, False, SLATE_500, PP_ALIGN.CENTER, line_spacing=9)
    else:
        add_rect(s, Inches(0.6), Inches(1.58), Inches(12.13), Inches(0.95), fill=SLATE_50, radius=True, line=SLATE_200)
        add_text_box(s, Inches(0.75), Inches(1.72), Inches(11.8), Inches(0.65), "MachineSystem  (HOAT_DONG / BAO_TRI / HONG)  →  MaintenancePlan  (kế hoạch năm, template-driven)  →  FaultRecord  (ghi lỗi, tra KB)  →  SparePart  (phụ tùng)", 8, True, NAVY, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_quality(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=GREEN)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(5), Inches(0.20), "QUALITY", 7, True, GREEN)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Chất lượng được ghi nhận theo mẻ — không phải cảm tính", 19, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.22), "Màu, mùi, vị, độ giòn — mỗi chỉ tiêu có người ghi, có mẻ gắn, có lịch sử", 8, False, SLATE_500, italic=True)
    for i,(t,d) in enumerate([("Auto-fill tỉ lệ","QualityEvaluation lấy tiLe tự động từ FinishedProduct — không nhập lại"),("Đánh giá đa chiều","Màu/mùi/vị/độ giòn, thang điểm chuẩn, có người ký & thời gian"),("Truy nguyên về mẻ","Mọi phiếu QC gắn maChien — lần ngược về mẻ chiên trong phút")] ):
        x=Inches(0.6)+i*Inches(4.15)
        add_rect(s, x, Inches(1.58), Inches(3.85), Inches(1.35), fill=WHITE, radius=True, line=SLATE_200)
        # icon circle
        ic=add_rect(s, x+Inches(0.20), Inches(1.72), Inches(0.36), Inches(0.36), fill=RGBColor(0xDC,0xFC,0xE7), radius=True)
        ic.text_frame.paragraphs[0].text=str(i+1); ic.text_frame.paragraphs[0].font.size=Pt(8); ic.text_frame.paragraphs[0].font.bold=True; ic.text_frame.paragraphs[0].font.color.rgb=GREEN; ic.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER; ic.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
        add_text_box(s, x+Inches(0.20), Inches(2.18), Inches(3.45), Inches(0.25), t, 8, True, NAVY)
        add_text_box(s, x+Inches(0.20), Inches(2.48), Inches(3.45), Inches(0.55), d, 7, False, SLATE_500, line_spacing=10)
    add_rect(s, Inches(0.6), Inches(3.35), Inches(12.13), Inches(0.48), fill=RGBColor(0xF0,0xFD,0xF4), radius=True, line=RGBColor(0xBB,0xF7,0xD0))
    add_text_box(s, Inches(0.75), Inches(3.48), Inches(11.8), Inches(0.28), "Chất lượng không phải lời nói — là phiếu QC gắn mẻ, có điểm số, có người ký, truy nguyên được.", 7.5, True, GREEN, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_ai(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(4), Inches(0.20), "AI ASSISTANT", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "AI hiểu ngữ cảnh xưởng — không phải chatbot chung chung", 18, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.22), "Hỏi bằng tiếng Việt tự nhiên — trả lời bằng quy trình thực, thao tác thực", 8, False, SLATE_500, italic=True)
    ai_img=os.path.join(ASSETS,"ai-pipeline.png")
    if os.path.exists(ai_img):
        add_image(s, ai_img, Inches(0.50), Inches(1.50), width=Inches(12.33))
    # agent line
    add_rect(s, Inches(0.6), Inches(2.55), Inches(12.13), Inches(0.38), fill=RGBColor(0xEF,0xF6,0xFF), radius=True, line=RGBColor(0xBF,0xDB,0xFE))
    add_text_box(s, Inches(0.70), Inches(2.62), Inches(11.9), Inches(0.24), "+ ReAct Agent  •  72 tools  •  department RBAC  •  max 5 iterations  —  thao tác thực tế qua ERP API  •  write cần user confirm", 6.5, True, NAVY, PP_ALIGN.CENTER)
    # mock chat
    mock=os.path.join(ASSETS,"mock-chat.png")
    if os.path.exists(mock):
        add_image(s, mock, Inches(0.6), Inches(3.10), width=Inches(6.0))
        add_text_box(s, Inches(0.6), Inches(5.38), Inches(6.0), Inches(0.18), "▲  Mock UI — ChatWidget (tiếng Việt, grounded + sources, agent confirm)", 6, False, SLATE_400, PP_ALIGN.CENTER)
        # right proofs
        add_rect(s, Inches(6.85), Inches(3.10), Inches(0.02), Inches(2.28), fill=SLATE_200)
        for i,(t,d) in enumerate([("Hybrid retrieval","Dense + Sparse → RRF k=60 → FlashRank\nKhông chỉ vector search đơn thuần"),("Grounded + Faithfulness","Chỉ trả lời từ CONTEXT\nLLM-as-judge giảm hallucination"),("72-tool agent + guard","Department RBAC, max 5 iter\nWrite cần confirm")] ):
            y=Inches(3.10)+i*Inches(0.82)
            add_rect(s, Inches(7.05), y, Inches(5.6), Inches(0.70), fill=SLATE_50, radius=True, line=SLATE_200)
            add_text_box(s, Inches(7.20), y+Inches(0.08), Inches(5.3), Inches(0.20), t, 7.5, True, NAVY)
            add_text_box(s, Inches(7.20), y+Inches(0.30), Inches(5.3), Inches(0.32), d, 6.5, False, SLATE_500, line_spacing=9)
    else:
        for i,(t,d) in enumerate([("Hybrid retrieval","Dense + Sparse → RRF k=60 → FlashRank"),("Grounded + Faithfulness","CONTEXT + LLM-as-judge"),("72-tool agent + guard","RBAC, max 5 iter, write confirm")] ):
            x=Inches(0.6)+i*Inches(4.15); add_rect(s, x, Inches(3.10), Inches(3.85), Inches(1.0), fill=SLATE_50, radius=True, line=SLATE_200)
            add_text_box(s, x+Inches(0.2), Inches(3.20), Inches(3.45), Inches(0.22), t, 7.5, True, NAVY)
            add_text_box(s, x+Inches(0.2), Inches(3.45), Inches(3.45), Inches(0.5), d, 6.5, False, SLATE_500)
    add_logo_footer(s); return s

def slide_face(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.32), Inches(0.6), Inches(0.05), fill=CYAN)
    add_text_box(s, Inches(0.6), Inches(0.47), Inches(7), Inches(0.20), "FACE KIOSK  •  PHYSICAL-DIGITAL BRIDGE", 7, True, CYAN)
    add_text_box(s, Inches(0.6), Inches(0.71), Inches(12), Inches(0.42), "Khuôn mặt là thẻ chấm công — 4 lớp chống giả mạo, học dần theo thời gian", 16, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.13), Inches(12), Inches(0.20), "Một tablet tại xưởng — không cần thẻ, không chấm hộ, không gian lận", 8, False, SLATE_500, italic=True)
    face=os.path.join(ASSETS,"face-layers.png")
    if os.path.exists(face):
        add_image(s, face, Inches(0.50), Inches(1.42), width=Inches(12.33))
    for i,t in enumerate(["Advisory lock — không duplicate","Adaptive gallery — càng dùng càng chính xác","Một tablet cho 100–300 CN — chi phí thấp"]):
        pill(s, Inches(0.6+i*4.15), Inches(3.58), t, bg=TEAL if i==0 else SLATE_100, fg=WHITE if i==0 else SLATE_700, font_size=6.5)
    # pipeline line
    add_rect(s, Inches(0.6), Inches(4.10), Inches(12.13), Inches(0.48), fill=RGBColor(0xEC,0xFE,0xFF), radius=True, line=RGBColor(0xA5,0xF3,0xFC))
    add_text_box(s, Inches(0.70), Inches(4.18), Inches(11.9), Inches(0.32), "RetinaFace / Yunet Detection  →  ArcFace Embedding (512D)  →  4-Layer Liveness  →  Voting Match  →  Attendance (advisory lock + dual cooldown)  →  Payroll / Timesheet  →  ERP", 6.5, True, NAVY, PP_ALIGN.CENTER, line_spacing=9)
    add_text_box(s, Inches(0.6), Inches(4.72), Inches(12.13), Inches(0.22), "Kiosk giá rẻ cho 100–300 CN  •  chưa PAD ISO 30107, không depth/IR  •  không claim FaceID phone-level — trung thực là điểm mạnh", 6.5, False, SLATE_400, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_tech(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(6), Inches(0.20), "TECHNOLOGY  •  5 LAYERS", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Không phải danh sách logo — là 5 layers có chủ ý", 19, True, NAVY)
    tech=os.path.join(ASSETS,"tech.png")
    if os.path.exists(tech):
        add_image(s, tech, Inches(0.55), Inches(1.32), width=Inches(7.0))
        # right note
        add_rect(s, Inches(7.85), Inches(1.32), Inches(5.0), Inches(4.1), fill=SLATE_50, radius=True, line=SLATE_200)
        add_text_box(s, Inches(8.05), Inches(1.48), Inches(4.6), Inches(0.28), "TẠI SAO 5 LAYERS NÀY?", 7, True, BLUE)
        add_text_box(s, Inches(8.05), Inches(1.82), Inches(4.6), Inches(3.3), "• Experience tách khỏi Application — SPA 46 pages độc lập\n• Data 3 schemas CUID + transaction + advisory locks\n• Intelligence tách riêng FastAPI — không lẫn business logic\n• Infrastructure Docker Compose + Nginx — single-VPS production\n• Mỗi layer một vai trò rõ — không thừa, không thiếu", 7, False, SLATE_700, line_spacing=11)
    else:
        y=Inches(1.32)
        for name,stack,role,col in [("EXPERIENCE","React 18 + TS + Vite 5 + TailwindCSS","SPA 46 pages",(37,99,235)),("APPLICATION","Node 20 + Express 5 + Prisma 6","82 routes",(20,184,166)),("DATA","PostgreSQL 16 + Redis 7","Multi-schema",(245,158,11)),("INTELLIGENCE","FastAPI + DeepSeek + ChromaDB","RAG + 72 tools",(6,182,212)),("INFRASTRUCTURE","Docker Compose + Nginx","Single-VPS",(51,65,85))]:
            add_rect(s, Inches(0.6), y, Inches(12.13), Inches(0.70), fill=WHITE, radius=True, line=SLATE_200)
            add_rect(s, Inches(0.6), y, Inches(0.05), Inches(0.70), fill=RGBColor(*col), radius=True)
            add_text_box(s, Inches(0.82), y+Inches(0.06), Inches(2), Inches(0.20), name, 6.5, True, RGBColor(*col))
            add_text_box(s, Inches(0.82), y+Inches(0.28), Inches(6), Inches(0.30), stack, 7, False, SLATE_700)
            y+=Inches(0.84)
    add_logo_footer(s); return s

def slide_arch(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(6), Inches(0.20), "ARCHITECTURE  •  MODULAR MONOLITH", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Modular monolith — đơn giản để deploy, đủ sâu để scale", 17, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.20), "Một Express app, 82 slices — deploy một lần, mở rộng từng module", 8, False, SLATE_500, italic=True)
    # infra diagram as styled card
    add_rect(s, Inches(0.6), Inches(1.52), Inches(12.13), Inches(1.85), fill=SLATE_50, radius=True, line=SLATE_200)
    # draw infra as text but with pill chips
    # top row: Internet → Nginx
    add_rect(s, Inches(0.80), Inches(1.72), Inches(1.5), Inches(0.38), fill=NAVY, radius=True)
    add_text_box(s, Inches(0.80), Inches(1.76), Inches(1.5), Inches(0.30), "Internet", 7, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.35), Inches(1.80), Inches(0.35), Inches(0.22), "→", 10, True, SLATE_400, PP_ALIGN.CENTER)
    add_rect(s, Inches(2.70), Inches(1.72), Inches(1.9), Inches(0.38), fill=BLUE, radius=True)
    add_text_box(s, Inches(2.70), Inches(1.76), Inches(1.9), Inches(0.30), "Nginx  TLS 1.2/1.3", 7, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(s, Inches(4.65), Inches(1.80), Inches(0.35), Inches(0.22), "→", 10, True, SLATE_400, PP_ALIGN.CENTER)
    # middle row: Frontend + Backend
    add_rect(s, Inches(0.80), Inches(2.22), Inches(1.9), Inches(0.42), fill=RGBColor(0xDB,0xE9,0xFE), radius=True, line=BLUE)
    add_text_box(s, Inches(0.80), Inches(2.26), Inches(1.9), Inches(0.34), "Frontend\nReact SPA  :5173", 6.5, True, NAVY, PP_ALIGN.CENTER, line_spacing=8)
    add_rect(s, Inches(2.90), Inches(2.22), Inches(2.1), Inches(0.42), fill=BLUE, radius=True)
    add_text_box(s, Inches(2.90), Inches(2.26), Inches(2.1), Inches(0.34), "Backend  Express 5  :5000", 6.5, True, WHITE, PP_ALIGN.CENTER, line_spacing=8)
    add_rect(s, Inches(5.20), Inches(2.22), Inches(1.8), Inches(0.42), fill=RGBColor(0x1E,0x3A,0x5F), radius=True)
    add_text_box(s, Inches(5.20), Inches(2.26), Inches(1.8), Inches(0.34), "PostgreSQL 16\n3 schemas", 6.5, True, WHITE, PP_ALIGN.CENTER, line_spacing=8)
    add_rect(s, Inches(7.20), Inches(2.22), Inches(1.5), Inches(0.42), fill=TEAL, radius=True)
    add_text_box(s, Inches(7.20), Inches(2.26), Inches(1.5), Inches(0.34), "Redis 7\ncache + RL", 6.5, True, WHITE, PP_ALIGN.CENTER, line_spacing=8)
    add_rect(s, Inches(8.90), Inches(2.22), Inches(2.2), Inches(0.42), fill=AMBER, radius=True)
    add_text_box(s, Inches(8.90), Inches(2.26), Inches(2.2), Inches(0.34), "AI Service  FastAPI :8001\nChromaDB + DeepSeek", 6.5, True, WHITE, PP_ALIGN.CENTER, line_spacing=8)
    # bottom row
    add_rect(s, Inches(0.80), Inches(2.78), Inches(2.6), Inches(0.38), fill=WHITE, radius=True, line=SLATE_200)
    add_text_box(s, Inches(0.80), Inches(2.84), Inches(2.6), Inches(0.26), "WebSocket  ws + Postgres NOTIFY", 6.5, True, SLATE_700, PP_ALIGN.CENTER)
    add_rect(s, Inches(3.60), Inches(2.78), Inches(2.2), Inches(0.38), fill=WHITE, radius=True, line=SLATE_200)
    add_text_box(s, Inches(3.60), Inches(2.84), Inches(2.2), Inches(0.26), "Web Push  VAPID", 6.5, True, SLATE_700, PP_ALIGN.CENTER)
    add_text_box(s, Inches(6.10), Inches(2.88), Inches(6), Inches(0.22), "cross-instance ready  •  advisory lock cho cron  •  sẵn sàng scale horizontal", 6.5, False, SLATE_400)
    for i,(t,d) in enumerate([("82 routes qua ROUTE_MAP","Đăng ký tập trung tại routes/index.ts\nKhông route nào silently ignored"),("Multi-schema Prisma","auth / business / common, CUID\nChild tables không JSON columns"),("Cross-instance ready","Postgres NOTIFY cho WS\nAdvisory lock cho cron — sẵn sàng scale")] ):
        x=Inches(0.6)+i*Inches(4.15)
        add_rect(s, x, Inches(3.65), Inches(3.85), Inches(1.15), fill=WHITE, radius=True, line=SLATE_200)
        add_text_box(s, x+Inches(0.20), Inches(3.75), Inches(3.45), Inches(0.22), t, 7.5, True, NAVY)
        add_text_box(s, x+Inches(0.20), Inches(4.05), Inches(3.45), Inches(0.50), d, 6.5, False, SLATE_500, line_spacing=9)
    add_logo_footer(s); return s

def slide_security(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=GREEN)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(7), Inches(0.20), "SECURITY  •  THỰC TẾ, KHÔNG KHẨU HIỆU", 7, True, GREEN)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Bảo vệ từ JWT tới kiện hàng — thực tế, không khẩu hiệu", 17, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.20), "Những gì đã có — và những gì chúng tôi không claim. Trung thực là cách tạo tin cậy.", 8, False, SLATE_500, italic=True)
    add_rect(s, Inches(0.6), Inches(1.52), Inches(5.95), Inches(0.32), fill=RGBColor(0xF0,0xFD,0xF4), radius=True, line=RGBColor(0xBB,0xF7,0xD0))
    add_text_box(s, Inches(0.75), Inches(1.58), Inches(5.6), Inches(0.20), "ĐÃ CÓ  (evidence)", 7, True, GREEN, PP_ALIGN.CENTER)
    add_rect(s, Inches(6.78), Inches(1.52), Inches(5.95), Inches(0.32), fill=AMBER_LIGHT, radius=True, line=RGBColor(0xFD,0xE6,0x8A))
    add_text_box(s, Inches(6.93), Inches(1.58), Inches(5.6), Inches(0.20), "CHƯA CLAIM  (roadmap)", 7, True, AMBER, PP_ALIGN.CENTER)
    left=["JWT access + refresh, role hierarchy ADMIN > HEAD > LEAD > EMP","RBAC/ABAC — secondary departments, 3 middlewares","Redis RL spoof-proof (IP + user)","Helmet, CORS, Zod validation","AES-GCM encrypt face embeddings","Audit log: face / evaluation / login","Advisory locks (attendance, face, cron)"]
    right=["2FA / MFA","WAF, vault secrets manager","Pentest / SOC2","mTLS / zero-trust","Backup offsite encrypted (hiện local /backup)","Audit toàn diện mọi entity","—"]
    add_rect(s, Inches(0.6), Inches(1.94), Inches(5.95), Inches(2.55), fill=WHITE, radius=True, line=SLATE_200)
    add_rect(s, Inches(6.78), Inches(1.94), Inches(5.95), Inches(2.55), fill=WHITE, radius=True, line=SLATE_200)
    for i,t in enumerate(left): add_text_box(s, Inches(0.80), Inches(2.04+i*0.33), Inches(5.55), Inches(0.25), "✓  "+t, 6.5, False, SLATE_700)
    for i,t in enumerate(right): add_text_box(s, Inches(6.98), Inches(2.04+i*0.33), Inches(5.55), Inches(0.25), "○  "+t, 6.5, False, SLATE_400)
    for i,t in enumerate(["RBAC/ABAC thực tế — secondary departments","RL không spoof được — không tin X-Forwarded-For mù quáng","Face embeddings mã hóa — AES-GCM, không plaintext"]):
        pill(s, Inches(0.6+i*4.15), Inches(4.78), t, bg=SLATE_900 if i==0 else SLATE_100, fg=WHITE if i==0 else SLATE_700, font_size=6.5)
    add_logo_footer(s); return s

def slide_roadmap(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.50), Inches(6), Inches(0.20), "ROADMAP  •  5 PHASES", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.74), Inches(12), Inches(0.45), "Từ Integrated ERP → Smart Factory — lộ trình có cơ sở kỹ thuật", 16, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.18), Inches(12), Inches(0.20), "Mỗi phase xây trên nền đã production-hardened — không hứa viễn vông", 8, False, SLATE_500, italic=True)
    phases=[("01","Integrated ERP","✅ Production","82 routes, 60+ entities, kiện/FIFO, maChien, HR/payroll, backup 3 lớp, 6-phase playbook",NAVY,BLUE_LIGHT),("02","AI-Assisted Ops","🔄 Đang hoàn thiện","RAG ≥100 QAs RAGAS, streaming faithfulness, fallback, p95 dashboards",AMBER,AMBER_LIGHT),("03","Data-Driven Factory","📋 Kế tiếp","Dashboard yield/Brix/OT, FAR/FRR face, audit toàn diện, S3 encrypted",TEAL,RGBColor(0xF0,0xFD,0xFA)),("04","Smart Factory","🔮 Roadmap","Multi-site, IoT telemetry, predictive maintenance (ML), mobile app",CYAN,RGBColor(0xEC,0xFE,0xFF)),("05","Enterprise Platform","🔮 Tầm nhìn","Multi-company, marketplace, supply chain tài chính, API ecosystem",SLATE_400,SLATE_50)]
    for i,(num,title,status,desc,col,bg) in enumerate(phases):
        x=Inches(0.6)+i*Inches(2.5)
        add_rect(s, x, Inches(1.62), Inches(2.32), Inches(2.95), fill=bg, radius=True, line=SLATE_200)
        add_rect(s, x, Inches(1.62), Inches(2.32), Inches(0.06), fill=col, radius=True)
        add_text_box(s, x+Inches(0.12), Inches(1.78), Inches(0.4), Inches(0.20), num, 6, True, col)
        add_text_box(s, x+Inches(0.52), Inches(1.78), Inches(1.6), Inches(0.20), status, 6, True, col)
        add_text_box(s, x+Inches(0.12), Inches(2.06), Inches(2.08), Inches(0.32), title, 8, True, NAVY)
        add_text_box(s, x+Inches(0.12), Inches(2.48), Inches(2.08), Inches(1.65), desc, 6.5, False, SLATE_500, line_spacing=9)
    add_text_box(s, Inches(0.6), Inches(4.92), Inches(12.13), Inches(0.22), "Hiện tại: single-VPS, single-site, single-company — đủ cho 1 nhà máy 100–300 CN  •  Multi-site / IoT là roadmap", 7, False, SLATE_400, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

def slide_cta(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, NAVY)
    add_rect(s, Inches(0.6), Inches(0.45), Inches(0.6), Inches(0.05), fill=AMBER)
    add_text_box(s, Inches(0.6), Inches(0.60), Inches(7), Inches(0.22), "GẶP CHÚNG TÔI TẠI GIAN HÀNG", 7, True, AMBER)
    add_text_box(s, Inches(0.6), Inches(0.88), Inches(8), Inches(0.55), "Xem kiosk và mẻ chiên trực tiếp", 26, True, WHITE)
    add_text_box(s, Inches(0.6), Inches(1.48), Inches(8), Inches(0.32), "Đừng chỉ đọc catalog — hãy chạm vào hệ thống đang chạy.", 9, False, BLUE_MID, italic=True)
    # QR placeholder card
    cols=[("Gặp trực tiếp","Gian hàng [số]\n[Tên triển lãm] — [Ngày]\nĐội ngũ An Bình Foods + kỹ thuật\ncó mặt cả ngày."),("QR CODE","Quét để xem demo\n[URL demo / landing]\n5 screenshots teaser:\nM.Evaluation  •  8 grades\nWarehouse kiện  •  ChatWidget\nDashboard1"),("Liên hệ","[Email]  •  [Hotline]\n[Website]\nĐặt lịch demo riêng —\nchúng tôi tới nhà máy khảo sát\nquy trình sấy.")]
    titles=["Gặp trực tiếp","Quét để xem thêm","Liên hệ"]
    for i,(body_title,body) in enumerate(cols):
        x=Inches(0.6)+i*Inches(2.65)
        if i==1:
            add_rect(s, x, Inches(2.08), Inches(2.35), Inches(2.85), fill=WHITE, radius=True)
            # QR box inside
            add_rect(s, x+Inches(0.45), Inches(2.22), Inches(1.45), Inches(1.45), fill=SLATE_900, radius=True)
            add_text_box(s, x+Inches(0.45), Inches(2.62), Inches(1.45), Inches(0.28), "QR", 16, True, WHITE, PP_ALIGN.CENTER)
            add_text_box(s, x+Inches(0.45), Inches(2.92), Inches(1.45), Inches(0.22), "[URL demo]", 6, False, WHITE, PP_ALIGN.CENTER)
            add_text_box(s, x+Inches(0.15), Inches(3.82), Inches(2.05), Inches(0.22), titles[i], 7, True, NAVY, PP_ALIGN.CENTER)
            add_text_box(s, x+Inches(0.15), Inches(4.10), Inches(2.05), Inches(0.70), body, 6.5, False, SLATE_500, PP_ALIGN.CENTER, line_spacing=9)
        else:
            add_rect(s, x, Inches(2.08), Inches(2.35), Inches(2.85), fill=WHITE, radius=True)
            add_text_box(s, x+Inches(0.15), Inches(2.28), Inches(2.05), Inches(0.22), titles[i], 8, True, NAVY, PP_ALIGN.CENTER)
            add_text_box(s, x+Inches(0.15), Inches(2.62), Inches(2.05), Inches(1.7), body, 7, False, SLATE_500, PP_ALIGN.CENTER, line_spacing=10)
    add_text_box(s, Inches(0.6), Inches(5.42), Inches(12.13), Inches(0.22), "Mỗi mẻ chiên đều có dữ liệu  •  Mỗi kiện đều có vị trí  •  Mỗi công nhân đều có khuôn mặt làm thẻ  •  AI hiểu xưởng của bạn", 7, False, BLUE_MID, PP_ALIGN.CENTER)
    # dashboard peek (small strip bottom)
    dash=os.path.join(ASSETS,"mock-dashboard.png")
    if os.path.exists(dash):
        add_image(s, dash, Inches(1.8), Inches(5.80), width=Inches(9.7), height=Inches(0.95))
        # rounded mask effect: add white rounded overlay frame
        add_rect(s, Inches(1.8), Inches(5.80), Inches(9.7), Inches(0.95), fill=None, line=WHITE, radius=True)
    return s

def slide_qa(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, WHITE)
    add_rect(s, Inches(0.6), Inches(0.38), Inches(0.6), Inches(0.05), fill=BLUE)
    add_text_box(s, Inches(0.6), Inches(0.53), Inches(5), Inches(0.22), "Q&A  •  DEMO TRỰC TIẾP", 7, True, BLUE)
    add_text_box(s, Inches(0.6), Inches(0.78), Inches(9), Inches(0.45), "Hỏi & Đáp — Demo trực tiếp", 26, True, NAVY)
    add_text_box(s, Inches(0.6), Inches(1.28), Inches(9), Inches(0.22), "3 câu hỏi gợi ý để phá băng — sau đó demo kiosk + ChatWidget + Dashboard", 8, False, SLATE_500, italic=True)
    qs=["“Brix và 8-grade yield được tính thế nào trong hệ thống?”","“Kiosk face hoạt động ở xưởng bụi / ánh sáng kém ra sao?”","“Triển khai cho nhà máy 200 người mất bao lâu?”"]
    for i,q in enumerate(qs):
        add_rect(s, Inches(0.6), Inches(1.72+i*0.62), Inches(7.5), Inches(0.48), fill=SLATE_50, radius=True, line=SLATE_200)
        add_text_box(s, Inches(0.85), Inches(1.84+i*0.62), Inches(7.1), Inches(0.28), f"0{i+1}  {q}", 8, False, SLATE_700)
    add_rect(s, Inches(8.55), Inches(1.72), Inches(4.15), Inches(1.86), fill=NAVY, radius=True)
    add_text_box(s, Inches(8.75), Inches(1.88), Inches(3.75), Inches(0.28), "DEMO TRỰC TIẾP  •  5 PHÚT", 7, True, AMBER, PP_ALIGN.CENTER)
    add_text_box(s, Inches(8.75), Inches(2.22), Inches(3.75), Inches(1.05), "①  Kiosk — chấm công trực tiếp (2’)\n②  ChatWidget — hỏi YCBG tiếng Việt (2’)\n③  Dashboard mẻ chiên 8-grade (1’)", 7, False, BLUE_MID, PP_ALIGN.CENTER, line_spacing=11)
    # dashboard preview
    dash=os.path.join(ASSETS,"mock-dashboard.png")
    if os.path.exists(dash):
        add_image(s, dash, Inches(0.60), Inches(3.82), width=Inches(12.10), height=Inches(1.20))
        add_rect(s, Inches(0.60), Inches(3.82), Inches(12.10), Inches(1.20), fill=None, line=SLATE_200, radius=True)
    add_text_box(s, Inches(0.6), Inches(5.32), Inches(12.13), Inches(0.28), "Mời anh/chị qua gian hàng ngay sau phần này — 2 phút là thấy hệ thống chạy thực tế.", 8, True, BLUE, PP_ALIGN.CENTER)
    add_logo_footer(s); return s

TRANSITIONS = [
    ("fade", {}),           # 01 cover — fade
    ("push", {"dir":"l"}),  # 02 agenda — push left
    ("fade", {}),           # 03 problem — fade
    ("push", {"dir":"l"}),  # 04 solution
    ("wipe", {"dir":"l"}),  # 05 connected
    ("fade", {}),           # 06 mfg1 — hero, keep calm
    ("push", {"dir":"l"}),  # 07 mfg2
    ("wipe", {"dir":"l"}),  # 08 warehouse
    ("fade", {}),           # 09 engineering
    ("fade", {}),           # 10 quality
    ("push", {"dir":"l"}),  # 11 AI
    ("wipe", {"dir":"l"}),  # 12 face
    ("fade", {}),           # 13 tech
    ("push", {"dir":"l"}),  # 14 arch
    ("fade", {}),           # 15 security
    ("push", {"dir":"l"}),  # 16 roadmap
    ("fade", {}),           # 17 CTA — grand fade
    ("fade", {}),           # 18 Q&A
]

def inject_timings(prs):
    for idx, slide in enumerate(prs.slides):
        dur=DURATIONS[idx] if idx<len(DURATIONS) else 10000
        el=slide._element
        NSMAP={"p":"http://schemas.openxmlformats.org/presentationml/2006/main"}
        # remove old transition if any
        old_trans = el.find("p:transition", NSMAP)
        if old_trans is not None:
            el.remove(old_trans)
        trans=etree.Element("{http://schemas.openxmlformats.org/presentationml/2006/main}transition",nsmap=NSMAP)
        cSld=el.find("p:cSld",NSMAP)
        el.insert(list(el).index(cSld)+1, trans) if cSld is not None else el.append(trans)
        trans.set("advOnClk","0"); trans.set("advTm",str(dur)); trans.set("spd","med")
        # varied transition type
        ttype, tattr = TRANSITIONS[idx] if idx < len(TRANSITIONS) else ("fade", {})
        if ttype == "fade":
            child = etree.SubElement(trans, "{http://schemas.openxmlformats.org/presentationml/2006/main}fade")
            child.set("thruBlk","1")
        elif ttype == "push":
            child = etree.SubElement(trans, "{http://schemas.openxmlformats.org/presentationml/2006/main}push")
            child.set("dir", tattr.get("dir","l"))
        elif ttype == "wipe":
            child = etree.SubElement(trans, "{http://schemas.openxmlformats.org/presentationml/2006/main}wipe")
            child.set("dir", tattr.get("dir","l"))
        elif ttype == "cover":
            child = etree.SubElement(trans, "{http://schemas.openxmlformats.org/presentationml/2006/main}cover")
            child.set("dir", tattr.get("dir","l"))
        timing=el.find("p:timing",NSMAP)
        if timing is None:
            timing=etree.Element("{http://schemas.openxmlformats.org/presentationml/2006/main}timing",nsmap=NSMAP); el.append(timing)
        if timing.find("p:tnLst",NSMAP) is None:
            tn=etree.SubElement(timing,"{http://schemas.openxmlformats.org/presentationml/2006/main}tnLst")
            par=etree.SubElement(tn,"{http://schemas.openxmlformats.org/presentationml/2006/main}par")
            ctn=etree.SubElement(par,"{http://schemas.openxmlformats.org/presentationml/2006/main}cTn",id="1",dur="indefinite",restart="never",nodeType="tmrPar")
            etree.SubElement(ctn,"{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst")
    pres=prs._element if hasattr(prs,"_element") else prs.element
    pr=pres.find("p:presentationPr",NS)
    if pr is None:
        pr=etree.Element("{http://schemas.openxmlformats.org/presentationml/2006/main}presentationPr",nsmap={"p":"http://schemas.openxmlformats.org/presentationml/2006/main"}); pres.insert(0,pr)
    show=pr.find("p:showPr",NS)
    if show is None: show=etree.SubElement(pr,"{http://schemas.openxmlformats.org/presentationml/2006/main}showPr")
    show.set("loop","1"); show.set("showType","kiosk")

def build():
    gen_assets()
    prs=Presentation(); prs.slide_width=W; prs.slide_height=H
    slide_cover(prs); slide_agenda(prs); slide_problem(prs); slide_solution(prs)
    slide_connected(prs); slide_mfg1(prs); slide_mfg2(prs); slide_warehouse(prs)
    slide_engineering(prs); slide_quality(prs); slide_ai(prs); slide_face(prs)
    slide_tech(prs); slide_arch(prs); slide_security(prs); slide_roadmap(prs)
    slide_cta(prs); slide_qa(prs)
    inject_timings(prs)
    prs.save(OUT)
    print(f"Saved {OUT} ({os.path.getsize(OUT)/1024:.1f} KB, {len(prs.slides)} slides)")
    for i,d in enumerate(DURATIONS): print(f"  {i+1:02d}: {d/1000:.0f}s")

if __name__=="__main__": build()
